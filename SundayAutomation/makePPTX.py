import os

from imagekitio import ImageKit
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


class PowerPointmaker:
    MAX_CHAR = 388

    def __init__(self, data, date):
        self.date = date

        self.makefirstReadingPPT(data["first_reading"], self.date)
        self.makePsalmPPT(data["psalm"], self.date)
        self.makeesecondReadingPPT(data["second_reading"], self.date)
        self.makeGospelPPT(data["gospel"], self.date)

    def makefirstReadingPPT(self, text, date):
        reading = "First Reading"
        self.newSlideHelper(text, date, reading)

    def makePsalmPPT(self, text, date):
        reading = "Psalm"
        self.newSlideHelper(text, date, reading, is_psalm=True)

    def makeesecondReadingPPT(self, text, date):
        reading = "Second Reading"
        self.newSlideHelper(text, date, reading)

    def makeGospelPPT(self, text, date):
        reading = "Gospel"
        self.newSlideHelper(text, date, reading)

    @staticmethod
    def newSlide(text, layout, prs):

        slide = prs.slides.add_slide(layout)
        slide.shapes.add_picture('bg.jpg', 0, 0, width=prs.slide_width)

        txBox = slide.shapes.add_textbox(0, 0, Inches(13), Inches(7))
        txBox.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.clear()

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.color.rgb = RGBColor(255, 255, 255)
        font = run.font
        font.name = 'Calibri'

        font.size = Pt(44)

    def newSlideHelper(self, text, date, reading_type, is_psalm=False):
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        prs.slide_width = 11887200
        prs.slide_height = 6686550

        char_list = list(text)
        split_text = []
        char_counter = 0
        slide_text = ''

        default_char_size = 395
        psalm_char_size = 270

        char_size = 0

        if is_psalm:
            char_size = psalm_char_size
        else:
            char_size = default_char_size

        for char in char_list:
            if char_counter >= char_size:
                split_text.append(slide_text)
                slide_text = ''
                char_counter = 0
            slide_text += char
            char_counter += 1

        if slide_text:  # Append the last chunk if it's not empty
            split_text.append(slide_text)

        print(split_text)

        for text_chunk in split_text:
            self.newSlide(text_chunk, blank_slide_layout, prs)

        if reading_type == "First Reading":
            prs.save(f"First Readings/{reading_type} - {date}.pptx")
        elif reading_type == "Psalm":
            prs.save(f"Psalms/{reading_type} - {date}.pptx")
        elif reading_type == "Second Reading":
            prs.save(f"Second Readings/{reading_type} - {date}.pptx")
        elif reading_type == "Gospel":
            prs.save(f"Gospels/{reading_type} - {date}.pptx")

